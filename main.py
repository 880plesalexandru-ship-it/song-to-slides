import os
import re

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Setează dimensiunile slide-ului la raportul 16:9
latime_slide = Inches(15.98)
inaltime_slide = Inches(8.98)

# Un dicționar global pentru a ține evidența numărului de apariții ale fiecărei cântări
# Cheia va fi titlul cântării, iar valoarea va fi numărul de apariții.
numar_aparitii = {}

def este_refren(paragraf):
    """
    Verifică dacă un paragraf este indentat (folosind tab sau riglă).
    Un vers indentat este considerat refren.
    """
    indent_prima_linie = paragraf.paragraph_format.first_line_indent
    indent_stanga = paragraf.paragraph_format.left_indent

    return (
        paragraf.text.startswith("\t") or
        (indent_prima_linie and indent_prima_linie > Inches(0)) or
        (indent_stanga and indent_stanga > Inches(0))
    )

def format_text(paragraf):
    """
    Formatează textul dintr-un paragraf Word, aplicând * pentru bold și _ pentru italic,
    dar numai dacă este necesar (ex. pentru versuri de tip 'strofa' sau 'refren').
    """
    text = ""
    for run in paragraf.runs:
        fragment = run.text
        if run.bold:
            fragment = f"*{fragment}*"
        if run.italic:
            fragment = f"_{fragment}_"
        text += fragment
    return text.rstrip()

def citeste_docx(fisier):
    """
    Citește cântările dintr-un document Word structurat și returnează o listă de cântări.
    Cu formatare adăugată (* pentru bold, _ pentru italic).
    """
    document = Document(fisier)
    lista_cantari = []
    cantare_curenta = {"numar": None, "titlu": None, "continut": [], "refrene": []}
    buffer_versuri = []
    buffer_refren = []
    tip_vers = None

    for paragraf in document.paragraphs:
        text = "".join(run.text for run in paragraf.runs).rstrip()
        if text is None:
            continue

        match = re.match(r"^(\d+)\.$", text)

        if match:
            if cantare_curenta["numar"]:
                if buffer_versuri:
                    if tip_vers != "refren":
                        cantare_curenta["continut"].append({"tip": tip_vers, "versuri": buffer_versuri})
                if buffer_refren:
                    cantare_curenta["refrene"].append(buffer_refren)
                lista_cantari.append(cantare_curenta)

            cantare_curenta = {"numar": int(match.group(1)), "titlu": None, "continut": [], "refrene": []}
            buffer_versuri = []
            buffer_refren = []
            tip_vers = "numar"
            
        # ... (restul funcției citeste_docx rămâne la fel)
        elif cantare_curenta["numar"] and not cantare_curenta["titlu"]:
            if este_refren(paragraf):
                text = text.lstrip("\t")
                tip_vers = "refren"
                if buffer_refren and not text:
                    cantare_curenta["refrene"].append(buffer_refren)
                    buffer_refren = []
                text = format_text(paragraf)
                buffer_refren.append(text)
            else:
                tip_vers = "strofa"
            cantare_curenta["titlu"] = text.lstrip("\t")
            text = format_text(paragraf)
            buffer_versuri.append(text)

        elif isinstance(text, str) and text:
            if este_refren(paragraf):
                tip_curent = "refren"
                text = text.lstrip("\t")
                if buffer_refren and not text:
                    cantare_curenta["refrene"].append(buffer_refren)
                    buffer_refren = []
                text = format_text(paragraf)
                buffer_refren.append(text)
            else:
                tip_curent = "strofa"
            
            if tip_vers and tip_curent != tip_vers:
                if tip_vers != "refren":
                    cantare_curenta["continut"].append({"tip": tip_vers, "versuri": buffer_versuri})
                buffer_versuri = []

            tip_vers = tip_curent
            if tip_curent == "strofa":
                text = format_text(paragraf)
                buffer_versuri.append(text)

        elif not isinstance(text, str):
            print("Eroare: Textul nu este un șir, ci:", type(text), text)

        else:
            if buffer_versuri:
                if tip_vers != "refren":
                    cantare_curenta["continut"].append({"tip": tip_vers, "versuri": buffer_versuri})
                buffer_versuri = []
                tip_vers = None
            if buffer_refren:
                cantare_curenta["refrene"].append(buffer_refren)
                buffer_refren = []
    
    # Salvează ultima cântare
    if cantare_curenta["numar"]:
        if buffer_versuri:
            if tip_vers != "refren":
                cantare_curenta["continut"].append({"tip": tip_vers, "versuri": buffer_versuri})
        if buffer_refren:
            cantare_curenta["refrene"].append(buffer_refren)
        lista_cantari.append(cantare_curenta)

    return lista_cantari

def parcurge_cantarile(lista_cantari):
    """
    Parcurge cantarile și le procesează.
    """
    for index, cantare in enumerate(lista_cantari):
        creaza_ppt(cantare)
        print(f"\rProcesat: {index + 1}/{len(lista_cantari)}", end="", flush=True)

def creaza_ppt(cantare):
    prezentare_ppt = Presentation()
    prezentare_ppt.slide_width = latime_slide
    prezentare_ppt.slide_height = inaltime_slide
    titlu_cantare = cantare["titlu"]
    titlu_cantare_curat = re.sub(r"^[^\w]+|[^\w]+$", "", titlu_cantare)
    titlu_ppt = f"{cantare['numar']:03} {titlu_cantare_curat}"
    titlu_cantare = titlu_cantare_curat
    continut = cantare["continut"]
    refrene = cantare.get("refrene", [])
    refren_final = cantare.get("refren_final", None)

    index_refren = 0
    titlu_adaugat = False

    # Aici verificăm și actualizăm numărul de apariții
    global numar_aparitii
    if cantare['numar'] in numar_aparitii:
        numar_aparitii[cantare['numar']] += 1
    else:
        numar_aparitii[cantare['numar']] = 1

    # Determinăm folderul de salvare în funcție de numărul de apariții
    nr_aparitii_curent = numar_aparitii[cantare['numar']]
    if nr_aparitii_curent == 1:
        folder_destinatie = "General"
    elif nr_aparitii_curent == 2:
        folder_destinatie = "Nunta"
    else:
        folder_destinatie = "Colinde"

    # ... (restul funcției creaza_ppt rămâne la fel)
    for index_slide, parte in enumerate(continut):
        if parte["tip"] == "strofa":
            slide = prezentare_ppt.slides.add_slide(prezentare_ppt.slide_layouts[5])
            fundal = slide.background
            umplere = fundal.fill
            umplere.solid()
            umplere.fore_color.rgb = RGBColor(0, 0, 0)
            if not titlu_adaugat:
                _creaza_titlu_slide(slide, titlu_cantare)
                titlu_adaugat = True
            _creaza_continut_slide(slide, parte, index_slide, len(continut))

            if refrene and index_refren < len(refrene):
                urmatorul_refren = refrene[index_refren]
                if urmatorul_refren and urmatorul_refren[0]:
                    text_footer = urmatorul_refren[0]
                else:
                    text_footer = "***"
                _creaza_footer_slide(slide, text_footer)

            if refrene:
                refrene_ramase = len(refrene) - index_refren
                slide = prezentare_ppt.slides.add_slide(prezentare_ppt.slide_layouts[5])
                fundal = slide.background
                umplere = fundal.fill
                umplere.solid()
                umplere.fore_color.rgb = RGBColor(0, 0, 0)
                if not titlu_adaugat:
                    _creaza_titlu_slide(slide, titlu_cantare)
                    titlu_adaugat = True
                _creaza_continut_refren_slide(slide, refrene[index_refren], index_refren, len(refrene))

                if index_refren + 1 < len(refrene) and refrene_ramase == len(continut) - index_slide:
                    index_refren += 1

                if index_slide + 1 < len(continut):
                    urmatoarea_strofa = continut[index_slide + 1]
                    if "versuri" in urmatoarea_strofa:
                        text_footer = urmatoarea_strofa["versuri"][0]
                    else:
                        text_footer = "***"
                    _creaza_footer_slide(slide, text_footer)

        elif parte["tip"] == "refren":
            slide = prezentare_ppt.slides.add_slide(prezentare_ppt.slide_layouts[5])
            fundal = slide.background
            umplere = fundal.fill
            umplere.solid()
            umplere.fore_color.rgb = RGBColor(0, 0, 0)
            if not titlu_adaugat:
                _creaza_titlu_slide(slide, titlu_cantare)
                titlu_adaugat = True
            _creaza_continut_refren_slide(slide, parte["versuri"], index_slide + 1, len(continut))

            if index_slide + 1 < len(continut):
                urmatoarea_strofa = continut[index_slide + 1]
                if "versuri" in urmatoarea_strofa:
                    text_footer = urmatoarea_strofa["versuri"][0]
                else:
                    text_footer = "***"
                _creaza_footer_slide(slide, text_footer)
                
    if refren_final:
        slide = prezentare_ppt.slides.add_slide(prezentare_ppt.slide_layouts[5])
        fundal = slide.background
        umplere = fundal.fill
        umplere.solid()
        umplere.fore_color.rgb = RGBColor(0, 0, 0)
        _creaza_continut_refren_slide(slide, refren_final, len(refrene), len(refrene))
        _creaza_footer_slide(slide, "***")

    # Salvăm prezentarea în folderul determinat
    cale_fisier = titlu_ppt
    cale_fisier = re.sub(r'[\\/:*?"<>|]', '', cale_fisier)
    cale_fisier = cale_fisier.rstrip('. ')
    cale_fisier = os.path.join(folder_destinatie, cale_fisier + ".pptx") # Aici este modificarea cheie
    prezentare_ppt.save(cale_fisier)

def _creaza_footer_slide(slide, text_footer):
    # ... (Funcția rămâne la fel)
    latime_footer = Inches(3)
    inaltime_footer = Inches(0.5)
    pozitie_stanga_footer = latime_slide - latime_footer - Inches(0.1)
    pozitie_sus_footer = inaltime_slide - inaltime_footer - Inches(0.1)

    footer = slide.shapes.add_textbox(pozitie_stanga_footer, pozitie_sus_footer, latime_footer, inaltime_footer)
    text_footer_curat = re.sub(r"^[^\w]+|[^\w]+$", "", text_footer)
    text_frame = footer.text_frame
    text_frame.text = text_footer_curat

    text_frame.clear()

    p = text_frame.add_paragraph()
    p.font.size = Pt(38)
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.RIGHT
    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    bold, italic = False, False
    fragment = ""
    for char in text_footer:
        if char == "*":
            if fragment:
                run = p.add_run()
                run.text = fragment
                run.font.bold = bold
                run.font.italic = italic
                fragment = ""
            bold = not bold
        elif char == "_":
            if fragment:
                run = p.add_run()
                run.text = fragment
                run.font.bold = bold
                run.font.italic = italic
                fragment = ""
            italic = not italic
        else:
            fragment += char

    if fragment:
        run = p.add_run()
        run.text = fragment
        run.font.bold = bold
        run.font.italic = italic

def _creaza_titlu_slide(slide, titlu_cantare):
    # ... (Funcția rămâne la fel)
    latime_titlu = Inches(3)
    inaltime_titlu = Inches(0.5)
    pozitie_stanga_titlu = latime_slide - latime_titlu - Inches(0.1)
    pozitie_sus_titlu = Inches(0.1)

    titlu = slide.shapes.add_textbox(pozitie_stanga_titlu, pozitie_sus_titlu, latime_titlu, inaltime_titlu)
    text_frame = titlu.text_frame
    text_frame.text = titlu_cantare

    p = text_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Calibri"
    p.font.color.rgb = RGBColor(255, 255, 255)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

def _creaza_continut_slide(slide, parte, index_slide, total_strofe):
    # ... (Funcția rămâne la fel)
    pozitie_stanga_text = Inches(0.1)
    pozitie_sus_text = Inches(0.1)
    latime_text = latime_slide - Inches(2)
    inaltime_text = Inches(1)
    numar_box = slide.shapes.add_textbox(pozitie_stanga_text, pozitie_sus_text, latime_text, inaltime_text)
    numar_frame = numar_box.text_frame
    numar_frame.clear()
    numar_frame.text = f"{index_slide + 1}/{total_strofe}"
    numar_frame.paragraphs[0].font.size = Pt(38)
    numar_frame.paragraphs[0].font.name = "Calibri"
    numar_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    numar_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    pozitie_stanga_versuri = Inches(0.1)
    pozitie_sus_versuri = Inches(0.5)
    latime_versuri = latime_slide - Inches(0.2)
    inaltime_versuri = inaltime_slide - Inches(1)

    content_box = slide.shapes.add_textbox(pozitie_stanga_versuri, pozitie_sus_versuri, latime_versuri, inaltime_versuri)
    text_frame = content_box.text_frame
    text_frame.clear()

    for vers in parte["versuri"]:
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Inches(0.1)

        while vers:
            match = re.search(r"(\*([^*]+)\*|_([^_]+)_)", vers)
            if match:
                start, end = match.span()
                if start > 0:
                    fragment = p.add_run()
                    fragment.text = vers[:start]

                fragment = p.add_run()
                fragment.text = match.group(2) or match.group(3)
                if match.group(1).startswith("*"):
                    fragment.font.bold = True
                if match.group(1).startswith("_"):
                    fragment.font.italic = True

                vers = vers[end:]
            else:
                fragment = p.add_run()
                fragment.text = vers
                break

        for run in p.runs:
            run.text = run.text
            run.font.size = Pt(48)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(255, 255, 255)

    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def _creaza_continut_refren_slide(slide, refren, index_slide, total_refrene):
    # ... (Funcția rămâne la fel)
    pozitie_stanga_text = Inches(0.1)
    pozitie_sus_text = Inches(0.1)
    latime_text = latime_slide - Inches(2)
    inaltime_text = Inches(1)
    numar_box = slide.shapes.add_textbox(pozitie_stanga_text, pozitie_sus_text, latime_text, inaltime_text)
    numar_frame = numar_box.text_frame
    numar_frame.clear()
    numar_frame.text = f"Refren {index_slide + 1}/{total_refrene}" if total_refrene > 1 else "Refren"
    numar_frame.paragraphs[0].font.size = Pt(38)
    numar_frame.paragraphs[0].font.name = "Calibri"
    numar_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    numar_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    pozitie_stanga_versuri = Inches(0.1)
    pozitie_sus_versuri = Inches(0.5)
    latime_versuri = latime_slide - Inches(0.2)
    inaltime_versuri = inaltime_slide - Inches(1.5)

    content_box = slide.shapes.add_textbox(pozitie_stanga_versuri, pozitie_sus_versuri, latime_versuri, inaltime_versuri)
    text_frame = content_box.text_frame
    text_frame.clear()

    for vers in refren:
        p = text_frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Inches(0.1)

        while vers:
            match = re.search(r"(\*([^*]+)\*|_([^_]+)_)", vers)
            if match:
                start, end = match.span()
                if start > 0:
                    fragment = p.add_run()
                    fragment.text = vers[:start]

                fragment = p.add_run()
                fragment.text = match.group(2) or match.group(3)
                if match.group(1).startswith("*"):
                    fragment.font.bold = True
                if match.group(1).startswith("_"):
                    fragment.font.italic = True

                vers = vers[end:]
            else:
                fragment = p.add_run()
                fragment.text = vers
                break

        for run in p.runs:
            run.text = run.text
            run.font.size = Pt(48)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(255, 255, 255)

    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Main Script
if __name__ == "__main__":
    fisier_docx = "caiet.docx"
    
    # Creăm folderele General, Nunta și Colinde dacă nu există deja
    if not os.path.exists("General"):
        os.makedirs("General")
    if not os.path.exists("Nunta"):
        os.makedirs("Nunta")
    if not os.path.exists("Colinde"):
        os.makedirs("Colinde")
    
    lista_cantari = citeste_docx(fisier_docx)
    parcurge_cantarile(lista_cantari)