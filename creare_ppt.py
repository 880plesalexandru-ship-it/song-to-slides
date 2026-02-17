import tkinter as tk
from tkinter import filedialog, messagebox

from docx import Document
from pptx import Presentation
from pptx.util import Inches


def load_docx():
    file_path = filedialog.askopenfilename(filetypes=[("Word Documents", "*.docx")])
    if file_path:
        doc = Document(file_path)
        content = "\n".join([para.text for para in doc.paragraphs])
        text_box.delete("1.0", tk.END)
        text_box.insert(tk.END, content)


def save_pptx():
    text = text_box.get("1.0", tk.END).strip()
    if not text:
        messagebox.showwarning("Avertisment", "Introduceți text înainte de a salva.")
        return

    slides = text.split("\n\n")
    ppt = Presentation()

    for slide_text in slides:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Slide gol
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = textbox.text_frame
        text_frame.text = slide_text

    file_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                             filetypes=[("PowerPoint Presentation", "*.pptx")])
    if file_path:
        ppt.save(file_path)
        messagebox.showinfo("Succes", "Fișierul PPTX a fost salvat cu succes!")


# Creare UI
root = tk.Tk()
root.title("Conversie DOCX/PPTX")
root.geometry("600x400")

frame = tk.Frame(root)
frame.pack(pady=10)

tk.Button(frame, text="Încarcă DOCX", command=load_docx).pack(side=tk.LEFT, padx=5)
tk.Button(frame, text="Salvează ca PPTX", command=save_pptx).pack(side=tk.RIGHT, padx=5)

text_box = tk.Text(root, wrap=tk.WORD, height=15, width=70)
text_box.pack(pady=10)

root.mainloop()
